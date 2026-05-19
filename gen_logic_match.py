#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
逻辑结构匹配模式 v2：
1. 从catalog读取模板页逻辑结构
2. 按同样逻辑结构整理OpenClaw内容
3. 通过坐标匹配找到每个槽的运行时引用
4. 蒸馏到目标字数 → 替换
"""

from pptx import Presentation
from pptx.oxml.ns import qn
import shutil, json

import os
_SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATE = os.path.join(_SKILL_DIR, "template", "hw_template.pptx")

def replace_text_keep_font(para, new_text):
    for run in para.runs:
        run.text = ""
    if para.runs:
        para.runs[0].text = new_text

def distill(target_len, text):
    """智能蒸馏到目标字数。
    - 短槽(≤10字，标题/标签)：激进压缩，精准匹配
    - 中槽(11-20字，摘要)：适度压缩
    - 长槽(>20字，正文)：优先保内容，最小压缩"""
    if len(text) <= target_len:
        return text

    if target_len <= 10:
        compress_rules = [
            ("协作调度", "调度"), ("长期记忆", "记忆"), ("工作空间", "空间"),
            ("角色人格", "人格"), ("用户画像", "画像"), ("浏览器引擎", "浏览"),
            ("多端接入", "接入"), ("文档处理", "文档"), ("事件总线", "事件"),
            ("消息队列", "MQ"), ("异步任务", "异步"), ("流式处理", "流式"),
            ("安全沙箱", "沙箱"), ("权限隔离", "隔离"), ("审计追踪", "审计"),
            ("数据加密", "加密"), ("配置驱动", "配置"), ("声明式部署", "部署"),
            ("热加载", "热载"), ("模板化", "模板"), ("会话上下文", "上下文"),
            ("工具调用", "调用"), ("记忆持久", "持久"), ("可扩展性", "可扩"),
            ("去中心化", "分布"), ("异步非阻塞", "异步"), ("声明式配置", "声明"),
            ("WhatsApp", "WA"), ("Channel", "通道"), ("Browser", "浏览"),
            ("Console", "终端"), ("Workspace", "空间"), ("MEMORY", "记忆"),
            ("PROFILE", "画像"), ("SOUL", "角色"),
        ]
        for a, b in compress_rules:
            if len(text) <= target_len:
                return text
            text = text.replace(a, b)
        for sep in ["、", "·", "；", "，", "。"]:
            if len(text) <= target_len:
                return text
            text = text.replace(sep, "")
        if len(text) > target_len:
            text = text[:target_len]
        return text.strip()

    # 正文模式：保内容优先
    while "  " in text and len(text) > target_len:
        text = text.replace("  ", " ")
    if len(text) <= int(target_len * 1.2):
        return text.strip()
    return text[:target_len].strip()

# ===== 加载catalog，选Slide 101 =====
with open("template_catalog_v2.json", "r", encoding="utf-8") as f:
    catalog = json.load(f)

SLIDE_NUM = 101
entry = catalog[SLIDE_NUM - 1]
slots = entry["text_slots"]

# 按(top, left)排序
slots.sort(key=lambda s: (s["top"], s["left"]))

# ===== 分析 Slide 101 的逻辑结构 =====
# ===== 精确角色分类（字号+粗体+字数+位置综合判断）=====
def classify_role(s):
    sz, bld, top, chars = s["font_size"], s["bold"], s["top"], s["char_count"]
    max_sz = max(x["font_size"] for x in slots)
    
    # 主标题：最大字号或大字号+顶部
    if sz >= 24 and top < 1.5:
        return "主标题"
    
    # 分节标题（架构层标注）：大字号+粗体
    if sz >= 16 and bld:
        return "分项标题"
    
    # 数据值：大字号但不粗也不在顶部
    if sz >= 24 and not bld and top >= 1.5:
        return "数据值"
    
    # 描述正文：中等字号+长文本 (>12字)
    if 12 <= sz < 20 and chars > 12:
        return "描述文"
    
    # 分类标签：中等字号+粗体+短文本
    if 12 <= sz < 16 and bld:
        return "分类标签"
    
    # 短标签：任何字号但字数少
    if 10 <= sz < 20 and chars <= 12:
        return "短标签"
    
    # 注释/来源：小字号
    if sz < 10:
        return "注释"
    
    return "正文"

for s in slots:
    s["role"] = classify_role(s)

print("=" * 60)
print(f"Slide {SLIDE_NUM} 逻辑结构分析")
print("=" * 60)
for i, s in enumerate(slots):
    print(f"  [{i:2d}] {s['role']:6s} {s['font_size']:3.0f}pt {s['char_count']:3d}字 "
          f"Y={s['top']:.1f}\" X={s['left']:.1f}\"  \"{s['text'][:45]}\"")

# 角色统计
from collections import Counter
rc = Counter(s["role"] for s in slots)
print(f"\n角色分布: {dict(rc)}")

# ===== 按 Slide 101 逻辑结构生成 OpenClaw 内容 =====
# 角色的数量和字数必须严格匹配
plan = []
for s in slots:
    plan.append({"role": s["role"], "len": s["char_count"], "new": "", "orig": s["text"],
                 "_top": s["top"], "_left": s["left"]})

# ===== 内容池：严格按字数分池，杜绝截断省略号 =====

# 1字池
oc_1 = ["端", "核", "存", "算", "网", "安", "基", "智"]

# 2字池 (主战场 — 需要25+项)
oc_2 = [
    "终端", "消息", "通道", "技能", "浏览", "文档",
    "调度", "编排", "存储", "画像", "引擎", "记忆",
    "安全", "合规", "插件", "工具", "协议", "路由",
    "缓存", "加密", "签名", "过滤", "拦截", "事件",
    "队列", "通知", "限流", "熔断", "降级", "重试",
    "幂等", "压缩", "解压", "回滚", "审计", "日志",
]

# 3字池
oc_3 = [
    "多智体", "工具链", "持久化", "可扩展", "自进化",
    "去中心", "异步化", "声明式",
]

# 5字池
oc_5 = [
    "插件化架构", "去中心调度", "声明式配置", "异步非阻塞",
    "多Agent框架",
]

# 7字池 (全中文)
oc_7 = [
    "容器化一键部署",        # 7字
    "多智能体调度层",        # 7字
    "插件化能力市场",        # 7字
    "记忆持久化引擎",        # 7字
    "多通道统一接入",        # 7字
]

# 10字池 (全中文)
oc_10 = [
    "拥抱智能打造协作平台",    # 10字
    "多智能体全栈解决方案",    # 10字
    "开放生态插件化架构",      # 10字
]

# 分项标题池 (2~13字，首5个是4层架构核心)
oc_items = [
    "交互通道", "通道", "插件", "引擎", "记忆",
    "智能体", "工具", "安全", "协议", "存储",
    "调度", "编排", "监控", "合规", "扩展",
]

# 描述文池 (长文本都用语义压缩，distill会处理)
oc_descs = [
    "终端多端接入 WA 通道推送",
    "Skills技能 浏览器 文档处理 检索",
    "Agent调度 上下文 工具调用编排",
    "角色人格 用户画像 记忆 空间",
    "事件总线 消息队列 异步任务",
    "配置驱动 声明式部署 热加载",
    "安全沙箱 权限隔离 审计追踪",
]

# 正文/分类标签池
oc_body = [
    "多智体", "工具链", "持久化", "可扩展", "自进化",
    "去中心化", "异步非阻塞", "声明式配置",
    "开放生态", "插件市场", "版本管理", "灰度发布",
    "多租户", "权限控制", "服务注册", "健康检查",
    "负载均衡", "链路追踪", "指标采集", "告警通知",
    "配置中心", "密钥管理", "证书轮换", "进程守护",
    "优雅关闭", "资源回收", "故障转移", "热更新",
    "回滚策略", "审批流程", "审计日志",
]

# 注释池
oc_notes = [
    "OpenClaw v1.0 多Agent协作框架",
    "2024 OpenClaw GDEClaw Team",
    "Confidential 内部使用",
    "版本 v1.0 编译 2024-Q4",
    "架构设计 虚空影虾 安总审核",
]

# ===== 智能池选择：根据目标字数自动选池 =====
def pick_from_pools(target_len, idx, pools_map):
    """根据目标字数从对应池取词"""
    best = None
    for width in sorted(pools_map.keys()):
        if width >= target_len:
            pool = pools_map[width]
            item = pool[idx % len(pool)]
            # 如果刚好等于target，完美
            if len(item) <= target_len:
                return distill(target_len, item)
            best = best or item
    # 从最小匹配池取，用distill压缩
    return distill(target_len, best or "—")

# 按字数组织池
LEN_POOLS = {1: oc_1, 2: oc_2, 3: oc_3, 5: oc_5, 7: oc_7, 10: oc_10}

# ===== 分配 =====
# 主标题
for p in plan:
    if p["role"] == "主标题" and not p["new"]:
        p["new"] = distill(p["len"], "OpenClaw多Agent协作架构")
        break

# 分项标题
item_idx = 0
for p in plan:
    if p["role"] == "分项标题" and not p["new"]:
        p["new"] = distill(p["len"], oc_items[item_idx % len(oc_items)])
        item_idx += 1

# 描述文
desc_idx = 0
for p in plan:
    if p["role"] == "描述文" and not p["new"]:
        p["new"] = distill(p["len"], oc_descs[desc_idx % len(oc_descs)])
        desc_idx += 1

# 短标签 — 用智能池
tag_counters = {}
for p in plan:
    w = p["len"]
    if p["role"] == "短标签" and not p["new"]:
        tag_counters[w] = tag_counters.get(w, 0)
        p["new"] = pick_from_pools(w, tag_counters[w], LEN_POOLS)
        tag_counters[w] += 1

# 正文/分类标签
body_idx = 0
for p in plan:
    if p["role"] in ("正文", "分类标签") and not p["new"]:
        p["new"] = distill(p["len"], oc_body[body_idx % len(oc_body)])
        body_idx += 1

# 注释
note_idx = 0
for p in plan:
    if p["role"] == "注释" and not p["new"]:
        p["new"] = distill(p["len"], oc_notes[note_idx % len(oc_notes)])
        note_idx += 1

# 数据值
dv_idx = 0
oc_datavals = ["100%", "∞", "99.99%", "1M+", "24/7", "0ms"]
for p in plan:
    if p["role"] == "数据值" and not p["new"]:
        p["new"] = distill(p["len"], oc_datavals[dv_idx % len(oc_datavals)])
        dv_idx += 1

# 终极兜底 — 绝不清空
fallback_idx = 0
for p in plan:
    if not p["new"]:
        p["new"] = distill(p["len"], oc_2[fallback_idx % len(oc_2)])
        p["_fallback"] = True
        fallback_idx += 1

print(f"\n内容分配结果:")
print(f"  主标题: 1个 → 已分配")
print(f"  分项标题: {sum(1 for p in plan if p['role']=='分项标题')}个 → {item_idx}个")
print(f"  描述文: {sum(1 for p in plan if p['role']=='描述文')}个 → {desc_idx}个")
print(f"  短标签: {sum(1 for p in plan if p['role']=='短标签')}个 → {sum(tag_counters.values())}个(智能池)")
print(f"  正文+分类标签: {sum(1 for p in plan if p['role'] in ('正文','分类标签'))}个 → {body_idx}个")
print(f"  注释: {sum(1 for p in plan if p['role']=='注释')}个 → {note_idx}个")
print(f"  数据值: {sum(1 for p in plan if p['role']=='数据值')}个 → {dv_idx}个")
print(f"  兜底填充: {sum(1 for p in plan if p.get('_fallback'))}个")
print(f"  ❌ 空槽: {sum(1 for p in plan if not p['new'])}个")

# ===== 执行替换 =====
OUTPUT = "openclaw_logic.pptx"
shutil.copy2(TEMPLATE, OUTPUT)
prs = Presentation(OUTPUT)

# 只保留目标slide
sldIdLst = prs.slides._sldIdLst
all_ids = list(sldIdLst)
target_rId = all_ids[SLIDE_NUM - 1].get(qn('r:id'))
for sld_id in reversed(all_ids):
    if sld_id.get(qn('r:id')) != target_rId:
        sldIdLst.remove(sld_id)

slide = prs.slides[0]

# 用坐标匹配：对每个shape的每个paragraph，找plan中最近的槽
replace_count = 0
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    left = shape.left/914400 if shape.left else 0
    top = shape.top/914400 if shape.top else 0
    
    for pi, para in enumerate(shape.text_frame.paragraphs):
        text = para.text.strip()
        if not text:
            continue
        
        # 在plan中找位置最接近的未使用槽
        best = None
        best_dist = 999
        for p in plan:
            if p.get("_matched"):
                continue
            dist = abs(p.get("_top", 0) - top) + abs(p.get("_left", 0) - left)
            if dist < best_dist:
                best_dist = dist
                best = p
        
        if best:
            replace_text_keep_font(para, best["new"])
            best["_matched"] = True
            print(f"  [{best['role']}] {best['orig'][:25]} → {best['new'][:35]} ({best['len']}→{len(best['new'])}字)")
            replace_count += 1
        else:
            # 无匹配也不清空 — 用最近兜底
            replace_text_keep_font(para, "—")
            print(f"  [?] {text[:25]} → — (无匹配兜底)")

prs.save(OUTPUT)
print(f"\n✓ {OUTPUT} ({replace_count}处替换, 格式/布局/字号完全不变)")
print("=" * 60)

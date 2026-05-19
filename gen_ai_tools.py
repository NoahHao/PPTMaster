#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""区域AI工具建设建议与治理规范 — 华为胶片大师 8页生成"""
import os, re, zipfile
from xml.etree import ElementTree as ET
from pptx import Presentation
from pptx.util import Emu

_SKILL_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = os.path.join(_SKILL_DIR, "template", "hw_template.pptx")
OUTPUT = os.path.join(os.path.dirname(_SKILL_DIR), "区域AI工具建设建议与治理规范.pptx")

NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
for p,u in [('p',NS_P),('r',NS_R)]: ET.register_namespace(p,u)
ET.register_namespace('',NS_P)

# Page assignment (source slide -> template slide)
# Src1:封面 → T2, Src2:目录 → T63, Src3:知识工程 → T54
# Src4:能力域 → T5, Src5:规范标准 → T37, Src6:合规流程 → T82
# Src7:安全红线 → T141, Src8:结束 → T3
PAGES = [2, 63, 54, 5, 37, 82, 141, 3]

# === Extract selected slides ===
with zipfile.ZipFile(TEMPLATE, 'r') as zin:
    rels = ET.fromstring(zin.read('ppt/_rels/presentation.xml.rels'))
    s2r = {}; r2s = {}
    for c in rels:
        m = re.search(r'slide(\d+)', c.get('Target',''))
        if m: s2r[int(m.group(1))] = c.get('Id'); r2s[c.get('Id')] = int(m.group(1))
    
    krids = {s2r[sn] for sn in PAGES}
    drids = {r for r in r2s if r not in krids}
    
    delf = set()
    for fn in zin.namelist():
        m = re.match(r'ppt/slides/slide(\d+)\.xml$', fn)
        if m and int(m.group(1)) not in PAGES: delf.add(fn)
        m = re.match(r'ppt/slides/_rels/slide(\d+)\.xml\.rels$', fn)
        if m and int(m.group(1)) not in PAGES: delf.add(fn)
    
    with zipfile.ZipFile(OUTPUT, 'w', zipfile.ZIP_DEFLATED) as zo:
        for item in zin.infolist():
            fn = item.filename
            if fn in delf: continue
            data = zin.read(fn)
            if fn == '[Content_Types].xml':
                r = ET.fromstring(data)
                for c in list(r):
                    if any('/'+df in (c.get('PartName','')) for df in delf):
                        r.remove(c)
                data = ET.tostring(r, encoding='UTF-8', xml_declaration=True)
            elif fn == 'ppt/presentation.xml':
                r = ET.fromstring(data)
                sl = r.find('.//{'+NS_P+'}sldIdLst')
                if sl is not None:
                    for c in list(sl):
                        if c.get('{'+NS_R+'}id') in drids: sl.remove(c)
                    rem = list(sl); sl.clear()
                    for sn in PAGES:
                        for c in rem:
                            if c.get('{'+NS_R+'}id') == s2r[sn]: sl.append(c); break
                data = ET.tostring(r, encoding='UTF-8', xml_declaration=True)
            elif fn == 'ppt/_rels/presentation.xml.rels':
                r = ET.fromstring(data)
                for c in list(r):
                    if c.get('Id') in drids: r.remove(c)
                data = ET.tostring(r, encoding='UTF-8', xml_declaration=True)
            zo.writestr(item, data)

print(f"✅ Extracted {len(PAGES)} template pages")

# === Fill engine ===
prs = Presentation(OUTPUT)

def all_text(slide):
    items = []
    def walk(shape):
        try:
            if hasattr(shape,'shapes'):
                for s in shape.shapes: walk(s)
        except: pass
        try:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    t = p.text.strip()
                    if t: items.append((shape.top or Emu(0), shape.left or Emu(0), shape, p, t))
        except: pass
    for s in slide.shapes: walk(s)
    items.sort(key=lambda x: (x[0], x[1]))
    return items

def dist(n, t):
    if len(t) <= n: return t
    if n <= 10:
        for s in ["，","、","。","；"," ","·"]:
            if len(t) <= n: return t
            t = t.replace(s, "")
        return t[:n]
    t = t.replace("  "," ")
    if len(t) <= int(n*1.2): return t.strip()
    return t[:n].strip()

def fill_all(slide, content):
    items = all_text(slide)
    for _,_,_,p,_ in items:
        for r in p.runs: r.text = ""
    for i,(_,_,_,p,orig) in enumerate(items):
        if i < len(content) and content[i]:
            nw = dist(len(orig), content[i]) if orig else content[i]
            if p.runs: p.runs[0].text = nw
            elif nw: p.add_run().text = nw

# === PAGE 1: 封面 (Template S2 → 华为愿景页) ===
print("\n--- P1 封面 ---")
fill_all(prs.slides[0], [
    '区域AI工具建设建议与治理规范',
    '以AI为引擎，构建分层分类、服务业务的知识管理新范式',
    '驱动区域组织能力持续进化',
    '2026 区域数字化转型专项汇报',
    '人工智能治理篇',
    '知识工程体系化',
    '能力域运营与规范',
    '资源层治理与安全红线',
    '构建合规、高效、可进化的AI工具生态',
])

# === PAGE 2: 目录 (Template S63 → 公司战略页) ===
print("\n--- P2 目录 ---")
fill_all(prs.slides[1], [
    '目录 CONTENTS',
    '01 知识工程体系化',
    '依托员工助手，构建分层级知识体系',
    '实现知识从个人到组织的有序流动与复用',
    '02 能力域运营与规范',
    '建立个人与业务效率类技能的运营策略',
    '规范标准与推广机制，促进技能有序管理',
    '03 资源层治理与安全红线',
    '明确外部与内部资源的合规接入流程',
    '建立严格的安全治理体系与安全红线',
])

# === PAGE 3: 知识工程体系化 (Template S54 → 小微企业 3角色) ===
print("\n--- P3 知识工程三层体系 ---")
fill_all(prs.slides[2], [
    '知识工程：三层递进体系',
    '个人层：易沉淀易检索',
    '个人工作笔记、学习心得、项目经验总结',
    '团队层：易共享易协作',
    '团队项目文档、技术方案、标准作业流程SOP',
    '组织层：易汇聚易复用',
    '经过验证的最佳实践、培训材料、战略知识库',
    '实现知识从个人到组织的有序流动',
    '将分散知识转化为组织核心竞争力',
    '员工助手为统一入口',
    '分层分类管理确保知识不流失',
    '从"人找知识"到"知识找人"',
    '赋能每一位员工成为知识工作者',
])

# === PAGE 4: 能力域运营 (Template S5 → 研发投入 KPI页) ===
print("\n--- P4 能力域运营 ---")
fill_all(prs.slides[3], [
    '02 能力域运营与规范',
    '个人办公效率类技能',
    '策略：全面开放共享，鼓励全员参与AI技能开发与分享',
    '目标：沉淀文档处理、信息整理、知识总结等高频效率工具',
    '机制：月度效率之星评选、技能市集、定期迭代',
    '业务效率类技能',
    '策略：按核心业务场景分类管理，技能与业务紧密耦合',
    '目标：聚焦业务流关键痛点，实现降本增效与业务价值创造',
    '机制：Owner负责制，业务线资深专家担任Skill Owner',
    '两类技能相辅相成',
    '个人效率赋能组织效能',
    '业务效率驱动数字化转型',
    '共同构建AI驱动的工作新范式',
])

# === PAGE 5: 规范标准与推广 (Template S37 → 客户需求/客户价值 双栏) ===
print("\n--- P5 规范标准与推广 ---")
fill_all(prs.slides[4], [
    '规范标准与推广机制',
    '规范标准 Standards',
    '命名规范：功能名称_姓名_日期 统一格式',
    '描述规范：Markdown格式，含功能说明与实现逻辑',
    '审核规范：代码审核+安全扫描+合规检查',
    '推广机制 Promotion',
    '区域Skill空间：打造统一展示与交流平台',
    '月度评选：效率之星+创新之星双维度评选',
    '培训赋能：Workshop+办公达人认证体系',
    '以规范保质量，以推广促应用',
    '建立技能全生命周期管理闭环',
    '从创建到退役的完整治理链路',
    '让每一个技能都可发现、可复用、可信赖',
])

# === PAGE 6: 合规接入流程 (Template S82 → 自智网络产业页) ===
print("\n--- P6 合规接入流程 ---")
fill_all(prs.slides[5], [
    '03 资源层治理：合规接入流程',
    '申请阶段 APPLICATION',
    '填报《AI工具接入申请表》',
    '明确工具类型、数据范围与使用场景',
    '评审阶段 REVIEW',
    '安全合规团队进行技术评审',
    '数据隐私影响评估+第三方服务安全审查',
    '准入阶段 APPROVAL',
    '通过评审后发放接入凭证',
    '纳入周期性复查清单，定期审计',
    '三阶段层层把关',
    '确保每一个AI工具安全可控',
    '外部工具严审，内部工具规范',
    '构建可信AI工具生态的基石',
])

# === PAGE 7: 安全红线 (Template S141 → 网络安全页) ===
print("\n--- P7 安全红线 ---")
fill_all(prs.slides[6], [
    '安全红线：严禁触碰的三条底线',
    '红线一：严禁数据出境',
    '严禁任何AI工具将敏感数据传输至境外服务器',
    '所有数据处理必须在合规区域内完成',
    '红线二：严禁未授权采集',
    '严禁AI工具在未获用户明示同意情况下采集个人信息',
    '采集行为必须透明且可追溯',
    '红线三：严禁黑盒决策',
    '严禁使用不可解释的AI模型做关键业务决策',
    '所有AI输出必须有可审计的推理链路',
    '安全是AI工具建设的最高纲领',
    '触碰红线者一票否决',
    '治理不是阻碍创新',
    '而是为创新构建安全的轨道',
])

# === PAGE 8: 结束页 (Template S3 → 华为愿景页) ===
print("\n--- P8 结束页 ---")
fill_all(prs.slides[7], [
    '以AI为引擎，驱动组织能力持续进化',
    '构建分层分类、服务业务的知识管理新范式',
    '让每一位员工都拥有AI时代的数字助手',
    '知识工程体系化',
    '能力域运营与规范',
    '资源层治理与安全红线',
    '三层架构 · 一体推进 · 持续迭代',
    '感谢聆听',
    'Together We Build the AI-Powered Future',
])

# Clear notes
for s in prs.slides:
    try:
        if s.has_notes_slide: s.notes_slide.notes_text_frame.clear()
    except: pass

prs.save(OUTPUT)
print(f"\n✅ 区域AI工具建设建议与治理规范.pptx ({os.path.getsize(OUTPUT)/1024:.0f} KB)")
print(f"   共 {len(prs.slides)} 页")

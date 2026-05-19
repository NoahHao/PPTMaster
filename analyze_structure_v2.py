#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
精准模板分析 v2：
- 保留每页所有文本的位置/字号/字体/字数
- 标注每页适合的内容结构（总分/总分总/分总/并列/...）
- 输出完整的文本替换映射表
"""

from pptx import Presentation
from pptx.oxml.ns import qn
import json, os

_SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATE = os.path.join(_SKILL_DIR, "template", "hw_template.pptx")
prs = Presentation(TEMPLATE)

def get_slide_dna(slide):
    """
    提取一页的完整DNA：每个文本框的位置、字号、字体、原文字、字数
    """
    text_slots = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for pi, para in enumerate(shape.text_frame.paragraphs):
                text = para.text
                if not text.strip():
                    continue
                
                # 字号和字体取第一个run
                size_pt = 0
                font_name = ""
                bold = False
                italic = False
                color = ""
                if para.runs:
                    r = para.runs[0]
                    size_pt = r.font.size/12700 if r.font.size else 0
                    font_name = r.font.name if r.font.name else ""
                    bold = r.font.bold or False
                    italic = r.font.italic or False
                    try:
                        color = str(r.font.color.rgb) if r.font.color and r.font.color.rgb else ""
                    except:
                        color = "SCHEME"
                
                left = shape.left/914400 if shape.left else 0
                top = shape.top/914400 if shape.top else 0
                w = shape.width/914400 if shape.width else 0
                h = shape.height/914400 if shape.height else 0
                
                text_slots.append({
                    "shape_name": shape.name,
                    "para_index": pi,
                    "text": text.strip(),
                    "char_count": len(text.strip()),
                    "font_size": round(size_pt, 1),
                    "font_name": font_name,
                    "bold": bold,
                    "italic": italic,
                    "color": color,
                    "left": round(left, 2),
                    "top": round(top, 2),
                    "width": round(w, 2),
                    "height": round(h, 2),
                })
    return text_slots

def classify_structure_type(slots, slide_num):
    """
    根据文本槽位特征判断适合的内容结构
    
    返回: (结构类型, 置信度, 说明)
    """
    if not slots:
        return ("纯图页", 0.5, "无文本框")
    
    # 按字号分组
    big_title = [s for s in slots if s["font_size"] >= 24 and s["bold"]]
    medium_title = [s for s in slots if 16 <= s["font_size"] < 24]
    body_text = [s for s in slots if 10 <= s["font_size"] < 16]
    small_text = [s for s in slots if s["font_size"] < 10]
    
    total_slots = len(slots)
    total_chars = sum(s["char_count"] for s in slots)
    
    # 封面型：1个超大标题 + 极少量其他
    if len(big_title) == 1 and total_slots <= 3 and total_chars < 100:
        return ("封面/标题页", 0.9, "大标题+副标题/日期")
    
    # 章节分隔：1个大标题 + 无正文
    if len(big_title) == 1 and len(body_text) == 0 and total_chars < 60:
        return ("章节分隔", 0.9, "纯大标题(章节封面)")
    
    # KPI数据型：2-6个超大数字 + 对应标签
    if 2 <= len(big_title) <= 6 and total_slots <= 12 and total_chars < 300:
        return ("KPI数据展示", 0.8, "大数字+小标签(适合3-6个数据指标)")
    
    # 总分型：1个大标题 + 3-6个中等说明
    if len(big_title) >= 1 and 3 <= len(body_text) <= 8:
        if len(big_title) == 1:
            return ("总分结构", 0.8, "1主标题+多个分点说明")
    
    # 总分总型：标题 + 概述段 + 分点 + 总结
    if len(big_title) >= 1 and len(body_text) >= 4 and total_chars > 200:
        # 检查是否有明显的总结性短文本
        has_summary = any(len(s["text"]) < 30 and s["font_size"] >= 14 for s in slots[-2:])
        if has_summary:
            return ("总分总结构", 0.7, "标题+展开+总结")
        return ("总分结构", 0.7, "标题+多段展开")
    
    # 分总型：先分点 + 最后总结大标题
    if len(body_text) >= 3 and len(big_title) >= 1 and big_title[-1]["top"] > body_text[-1]["top"]:
        return ("分总结构", 0.7, "分点展开→总结标题")
    
    # 并列型：多个同等级标题/卡片
    if len(medium_title) >= 3 and len(big_title) <= 1:
        char_variance = max(s["char_count"] for s in medium_title) - min(s["char_count"] for s in medium_title)
        if char_variance < 30:
            return ("并列结构", 0.8, "多个同等卡片(适合并列项)")
    
    # 大图+标语
    if total_slots <= 4 and total_chars < 100 and len(big_title) >= 1:
        return ("大图+标语", 0.8, "大背景图+简短标题/标语")
    
    # 图表页
    if total_slots <= 5 and len(big_title) == 1 and total_chars < 80:
        return ("图表数据页", 0.6, "标题+图表(文字极少)")
    
    # 默认
    return ("通用内容", 0.5, f"{total_slots}个文本框/{total_chars}字的通用页")


print("=" * 70)
print("模板页结构标注")
print("=" * 70)

catalog = []

for i, slide in enumerate(prs.slides):
    slide_num = i + 1
    slots = get_slide_dna(slide)
    struct_type, confidence, desc = classify_structure_type(slots, slide_num)
    
    # 生成简要描述
    text_summary = []
    for s in slots[:5]:  # 前5个文本
        t = s["text"][:40]
        text_summary.append(f"[{s['font_size']:.0f}pt×{s['char_count']}字]{t}")
    
    entry = {
        "slide": slide_num,
        "struct_type": struct_type,
        "confidence": confidence,
        "description": desc,
        "slot_count": len(slots),
        "total_chars": sum(s["char_count"] for s in slots),
        "text_slots": slots,
    }
    catalog.append(entry)
    
    # 简要输出
    mark = "★" if confidence >= 0.8 else ("·" if confidence >= 0.6 else "?")
    slot_info = "; ".join(text_summary[:3])
    print(f"{mark} Slide {slide_num:3d} | {struct_type:12s} | {len(slots):2d}槽 {sum(s['char_count'] for s in slots):4d}字 | {slot_info[:100]}")

# 统计
from collections import Counter
print(f"\n{'=' * 70}")
print("结构类型分布")
print("=" * 70)
types = Counter(c["struct_type"] for c in catalog)
for t, count in types.most_common():
    print(f"  {t}: {count} 页")

# 保存
with open("template_catalog_v2.json", "w", encoding="utf-8") as f:
    json.dump(catalog, f, ensure_ascii=False, indent=2)

print(f"\n✓ 完整目录已保存: template_catalog_v2.json")
print(f"  包含 {len(catalog)} 页的完整文本槽位信息")

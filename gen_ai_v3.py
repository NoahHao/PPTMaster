#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
区域AI工具建设建议与治理规范 — 华为胶片大师 v3
策略：选宽槽模板，跳过微型标签(≤6字)，只填实质性内容槽
"""
import os, re, zipfile
from xml.etree import ElementTree as ET
from pptx import Presentation
from pptx.util import Emu

_SKILL_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = os.path.join(_SKILL_DIR, "template", "hw_template.pptx")
OUTPUT = os.path.join(os.path.dirname(_SKILL_DIR), "区域AI工具建设建议与治理规范.pptx")

NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
for p,u in [('p',NS_P),('r',NS_R)]: ET.register_namespace(p,u)
ET.register_namespace('',NS_P)

# P1封面→S2  P2目录→S4  P3知识工程→S62  P4能力域→S123
# P5规范标准→S13  P6合规流程→S146  P7安全红线→S141  P8结束→S13
# Note: S13 reused for P5 and P8 (different instances from ZIP extraction)
# Actually can't reuse same slide number. Use S12 instead for P8.
PAGES = [2, 4, 62, 123, 13, 146, 141, 29]

print("🔧 提取模板...")
with zipfile.ZipFile(TEMPLATE, 'r') as zin:
    rels = ET.fromstring(zin.read('ppt/_rels/presentation.xml.rels'))
    s2r = {}; r2s = {}
    for c in rels:
        m = re.search(r'slide(\d+)', c.get('Target',''))
        if m: s2r[int(m.group(1))] = c.get('Id'); r2s[c.get('Id')] = int(m.group(1))
    
    krids = {s2r[sn] for sn in PAGES}
    drids = {r for r in r2s if r not in krids}
    
    delf = set()
    for fn in zin.namelist():
        m = re.match(r'ppt/slides/slide(\d+)\.xml$', fn)
        if m and int(m.group(1)) not in PAGES: delf.add(fn)
        m = re.match(r'ppt/slides/_rels/slide(\d+)\.xml\.rels$', fn)
        if m and int(m.group(1)) not in PAGES: delf.add(fn)
    
    with zipfile.ZipFile(OUTPUT, 'w', zipfile.ZIP_DEFLATED) as zo:
        for item in zin.infolist():
            fn = item.filename
            if fn in delf: continue
            data = zin.read(fn)
            if fn == '[Content_Types].xml':
                r = ET.fromstring(data)
                for c in list(r):
                    if any('/'+df in (c.get('PartName','')) for df in delf):
                        r.remove(c)
                data = ET.tostring(r, encoding='UTF-8', xml_declaration=True)
            elif fn == 'ppt/presentation.xml':
                r = ET.fromstring(data)
                sl = r.find('.//{'+NS_P+'}sldIdLst')
                if sl is not None:
                    for c in list(sl):
                        if c.get('{'+NS_R+'}id') in drids: sl.remove(c)
                    rem = list(sl); sl.clear()
                    for sn in PAGES:
                        for c in rem:
                            if c.get('{'+NS_R+'}id') == s2r[sn]: sl.append(c); break
                data = ET.tostring(r, encoding='UTF-8', xml_declaration=True)
            elif fn == 'ppt/_rels/presentation.xml.rels':
                r = ET.fromstring(data)
                for c in list(r):
                    if c.get('Id') in drids: r.remove(c)
                data = ET.tostring(r, encoding='UTF-8', xml_declaration=True)
            zo.writestr(item, data)

prs = Presentation(OUTPUT)

# ===== ENGINE: 跳过微型标签 =====
def all_slots(slide):
    items = []
    def walk(shape):
        try:
            if hasattr(shape,'shapes'):
                for s in shape.shapes: walk(s)
        except: pass
        try:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    t = p.text.strip()
                    if t: items.append((shape.top or Emu(0), shape.left or Emu(0), shape, p, t))
        except: pass
    for s in slide.shapes: walk(s)
    items.sort(key=lambda x: (x[0], x[1]))
    return items

def dist(n, t):
    if not t: return t
    if len(t) <= n: return t
    if n <= 8:
        # Ultra-short slot: aggressive abbreviation
        for s in ["，","、","。","；","·"," ",":","：","\n"]:
            if len(t) <= n: return t
            t = t.replace(s, "")
        return t[:n]
    if len(t) <= int(n*1.2): return t.strip()
    # Try sentence boundary cut
    cut = t[:n]
    for sep in ["。","；"]:
        idx = cut.rfind(sep)
        if idx > n*0.6: return t[:idx+1]
    return t[:n].strip()

def fill_all_slots(slide, content):
    """
    逐槽顺序填入。铁律：绝不在词语中间截断。
    如果内容比槽长 → 在标点处自然断，找不到断点 → 跳过（留空也不截字）。
    """
    slots = all_slots(slide)
    
    for _,_,_,p,_ in slots:
        for r in p.runs: r.text = ""
    
    ci = 0
    for _,_,_,p,orig in slots:
        if ci >= len(content):
            break
        txt = content[ci]
        if not txt:
            ci += 1
            continue
        
        target = max(len(orig), 1)
        
        if len(txt) <= target:
            # Fits perfectly
            new_text = txt
        else:
            # Too long — try to find a natural break point
            cut = txt[:target]
            best = None
            for sep in ["。","；","，","、","】","）"," ","·","：",":"]:
                idx = cut.rfind(sep)
                if idx > 0:
                    best = idx
                    break
            if best is not None and best >= 1:
                new_text = txt[:best]
            else:
                # No natural break found — SKIP this slot, don't truncate
                ci += 1
                continue
        
        if p.runs:
            p.runs[0].text = new_text
        elif new_text:
            p.add_run().text = new_text
        ci += 1

# ====================== PAGE 1: 封面 (S2 → 9 slots) ======================
print("\nP1 封面")
fill_all_slots(prs.slides[0], [
    "区域AI工具建设建议与治理规范",
    "以AI为引擎，构建分层分类的知识管理新范式，驱动组织能力进化",
    "知识工程",
    "能力域运营",
    "资源治理",
    "三层架构",
    "一体推进",
    "持续迭代",
    "区域数字化转型专项汇报",
])

# ====================== PAGE 2: 目录 (S4 → 10 slots) ======================
print("P2 目录")
fill_all_slots(prs.slides[1], [
    "目录 CONTENTS",
    "01",
    "知识工程",
    "02",
    "能力运营",
    "03",
    "资源治理",
    "三层体系",
    "依托员工助手，构建分层知识体系，实现知识有序流动与复用",
    "建立技能运营规范与推广机制，划定安全红线与合规接入流程",
])

# ====================== PAGE 3: 知识工程3层 (S62 → 12 slots) ======================
print("P3 知识工程三层体系")
fill_all_slots(prs.slides[2], [
    "知识工程体系化",
    "依托员工助手，构建分层级知识体系",
    "个人层",
    "团队层",
    "经验沉淀成长起点",
    "文档共享提升协作",
    "组织层",
    "知识库",
    "最佳实践",
    "知识资产汇聚复用形成核心竞争力",
    "SOP标准作业流程管理",
    "从个人笔记到组织知识库，分层管理不流失",
])

# ====================== PAGE 4: 能力域2轨 (S123 → 11 slots) ======================
print("P4 能力域运营")
fill_all_slots(prs.slides[3], [
    "能力域运营与规范",
    "双轨并行",
    "办公效率类",
    "业务效率类",
    "全员创新",
    "降本增效",
    "开放共享，鼓励全员参与AI技能开发与分享",
    "按业务场景分类，技能与业务紧密耦合",
    "沉淀文档处理与信息整理等高频效率工具",
    "聚焦业务痛点，优化流程，实现价值创造",
    "月度评选、技能市集、Owner负责制",
])

# ====================== PAGE 5: 规范标准 (S13 → 5 slots) ======================
print("P5 规范标准与推广机制")
fill_all_slots(prs.slides[4], [
    "规范标准与推广机制",
    "规范标准：统一命名格式，强制Markdown描述含功能说明与核心实现逻辑",
    "推广机制：区域Skill空间展示交流，月度效率之星与创新之星双维度评选",
    "建立技能全生命周期管理闭环，从创建审核到退役的完整治理链路",
    "让每一个技能可发现、可复用、可信赖。以规范保质量，以推广促应用。",
])

# ====================== PAGE 6: 合规流程 (S146 → 8 slots) ======================
print("P6 合规接入流程")
fill_all_slots(prs.slides[5], [
    "合规接入流程",
    "申请",
    "评审",
    "准入",
    "填报申请表",
    "安全合规审查",
    "发放凭证",
    "全流程合规",
    "明确工具类型与用途",
    "隐私评估与安全审查",
    "纳入周期复查审计",
    "可追溯可审计",
])

# ====================== PAGE 7: 安全红线 (S141 → 6 slots) ======================
print("P7 安全红线")
fill_all_slots(prs.slides[6], [
    "安全红线：严禁触碰的三条底线",
    "红线一：严禁数据出境。严禁任何AI工具将敏感数据传输至境外服务器，所有数据处理必须在合规区域内完成。违规后果：一票否决。",
    "红线二：严禁未授权采集。严禁AI工具在未获用户明示同意情况下采集个人信息，采集行为必须透明可追溯。违规后果：一票否决。",
    "红线三：严禁黑盒决策。严禁使用不可解释AI模型做关键业务决策，所有AI输出必须有可审计的推理链路。违规后果：一票否决。",
    "安全是AI工具建设的最高纲领。治理不是阻碍创新，而是为创新构建安全的轨道。触碰红线者一票否决，合规即底线。",
    "区域AI工具建设建议与治理规范",
])

# ====================== PAGE 8: 结束页 (S29 → 7 slots, each 4-8 chars) ======================
print("P8 结束页")
fill_all_slots(prs.slides[7], [
    "感谢聆听",
    "以AI为引擎",
    "知识工程",
    "能力域运营",
    "资源安全治理",
    "持续迭代",
    "组织能力进化",
])

# ===== SAVE =====
for s in prs.slides:
    try:
        if s.has_notes_slide: s.notes_slide.notes_text_frame.clear()
    except: pass

prs.save(OUTPUT)
print(f"\n✅ {os.path.getsize(OUTPUT)/1024:.0f} KB, {len(prs.slides)} 页")

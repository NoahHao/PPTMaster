#!/usr/bin/env python3
"""
华为胶片大师 · 模板结构分类图鉴（2页正文版）
模板 S163（6结构类型卡片）+ S98（5决策路径）
匹配策略：文本内容 + 位置 定位shape ➔ 纯文字替换
"""
import sys, os, json, copy
from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn

sys.path.insert(0, os.path.join(os.path.dirname(__file__), 'scripts'))
from template_utils import replace_text_keep_font, distill

TEMPLATE_PPT = os.path.join(os.path.dirname(__file__), 'template', 'hw_template.pptx')
CATALOG_JSON = os.path.join(os.path.dirname(__file__), 'references', 'template_catalog_v2.json')
OUTPUT_PPT = os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', '..', '华为胶片大师_模板结构分类图鉴.pptx')

with open(CATALOG_JSON, 'r', encoding='utf-8') as f:
    catalog = json.load(f)

entry163 = next(e for e in catalog if e['slide'] == 163)
entry98 = next(e for e in catalog if e['slide'] == 98)

PPT_IDX_163 = 100
PPT_IDX_98 = 61

print("=" * 60)
print("华为胶片大师 · 模板结构分类图鉴")
print(f"  Page 1: S163 (PPT idx {PPT_IDX_163})")
print(f"  Page 2: S98  (PPT idx {PPT_IDX_98})")
print("=" * 60)


def find_shape_by_text(slide, target_text, fuzzy=False):
    """Find a shape in the slide whose text matches (or contains) target_text."""
    target = target_text.strip()
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        shape_text = shape.text_frame.text.strip()
        if fuzzy:
            if target in shape_text or shape_text in target:
                return shape
        else:
            if shape_text == target:
                return shape
    return None


def find_shape_by_pos_text(slide, cat_slot, fuzzy_threshold=0.8):
    """Find shape matching catalog slot by text similarity + position proximity."""
    target = cat_slot['text'].strip()
    t_top = cat_slot['top']
    t_left = cat_slot['left']
    
    best_shape = None
    best_score = 0
    
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        st = shape.text_frame.text.strip()
        if not st:
            continue
        
        # Text similarity
        text_score = 0
        if st == target:
            text_score = 1.0
        elif target in st or st in target:
            text_score = 0.9
        elif len(st) >= 3 and len(target) >= 3:
            # Check partial match
            common = set(st) & set(target)
            if len(common) >= min(len(st), len(target)) * 0.5:
                text_score = 0.5
        
        if text_score < 0.5:
            continue
        
        # Position score (normalized)
        try:
            s_top = shape.top / 914400  # EMU to inches
            s_left = shape.left / 914400
        except:
            continue
        
        pos_dist = abs(s_top - t_top) + abs(s_left - t_left)
        pos_score = max(0, 1 - pos_dist / 5.0)  # Normalize
        
        score = text_score * 0.7 + pos_score * 0.3
        if score > best_score:
            best_score = score
            best_shape = shape
    
    return best_shape


def replace_shape_text(shape, new_text):
    """Replace the first paragraph's text in a shape."""
    para = shape.text_frame.paragraphs[0]
    if para.text.strip():
        replace_text_keep_font(para, new_text)
        return True
    return False


# ============================================================
# Build catalog slot → content mappings
# ============================================================

# --- S163 mappings ---
slots163 = entry163['text_slots']

# Classify slots by role (based on font_size, bold, char_count)
title_slots = [s for s in slots163 if s['font_size'] >= 24]
bold_labels = [s for s in slots163 if s.get('bold') and 14 <= s['font_size'] < 28 and s['char_count'] >= 4]
# Sort BOLD labels by position
bold_labels.sort(key=lambda s: (s['top'], s['left']))

# Also pick the best description slots to pair with BOLD labels
# For S163, some desc slots are close to BOLD labels (same top±0.3")
desc_slots = [s for s in slots163 if not s.get('bold') and s['font_size'] <= 12 and 6 <= s['char_count'] <= 30]
desc_slots.sort(key=lambda s: (s['top'], s['left']))

# Structure type data
struct_types = [
    {"name": "并列结构", "count": "79页·50.3%", "tag": "主力·多主题并行"},
    {"name": "通用内容", "count": "67页·42.7%", "tag": "灵活·跨度最大"},
    {"name": "总分总",   "count": "5页·3.2%",   "tag": "骨架·汇报叙事"},
    {"name": "KPI数据",  "count": "3页·1.9%",   "tag": "指标·数据展示"},
    {"name": "总分结构", "count": "2页·1.3%",   "tag": "展开·产品方案"},
    {"name": "纯图页",   "count": "1页·0.6%",   "tag": "零槽·全图展示"},
]

# Build slot→content pairs (catalog_slot_index → new_text)
slot_replacements_163 = []

# Main title
if title_slots:
    slot_replacements_163.append((title_slots[0], "华为胶片大师 · 157页模板结构分类图鉴"))

# 6 BOLD labels → structure names
for i, bs in enumerate(bold_labels[:6]):
    if i < len(struct_types):
        name = struct_types[i]['name']
        slot_replacements_163.append((bs, name))

# Pair each BOLD label with a nearby description slot for count data
for i in range(min(6, len(struct_types))):
    bs = bold_labels[i]
    # Find the closest desc slot to this BOLD label (same vertical region)
    nearest_desc = None
    min_dist = 999
    for ds in desc_slots:
        # Check if already used
        if any(ds is s for s, _ in slot_replacements_163):
            continue
        v_dist = abs(ds['top'] - bs['top'])
        if v_dist < 1.5 and v_dist < min_dist:
            min_dist = v_dist
            nearest_desc = ds
    
    if nearest_desc:
        count_text = struct_types[i]['count']
        slot_replacements_163.append((nearest_desc, count_text))


# --- S98 mappings ---
slots98 = entry98['text_slots']

title98 = [s for s in slots98 if s['font_size'] >= 24]
bold98 = [s for s in slots98 if s.get('bold') and 14 <= s['font_size'] < 28 and s['char_count'] >= 3]
desc98 = [s for s in slots98 if not s.get('bold') and s['font_size'] <= 14 and s['char_count'] >= 10]

bold98.sort(key=lambda s: (s['top'], s['left']))
desc98.sort(key=lambda s: (s['top'], s['left']))

print(f"\nS98: {len(bold98)} BOLD, {len(desc98)} desc, {len(title98)} titles")

# Decision data
decisions = [
    {"scenario": "多主题并行", "path": "并列结构·79页", "rec": "S13/S141/S146宽槽"},
    {"scenario": "单主题深展", "path": "通用内容·67页", "rec": "S62/S123灵活适配"},
    {"scenario": "数据指标汇报", "path": "KPI数据·3页", "rec": "S7/S9/S10集中展示"},
    {"scenario": "完整叙事汇报", "path": "总分总·5页骨架", "rec": "S2封面→S4目录→S29"},
    {"scenario": "封面过渡结尾", "path": "通用微型·8页", "rec": "S2/S29/S253精炼"},
]

slot_replacements_98 = []

# Main title
if title98:
    biggest = max(title98, key=lambda s: s['font_size'])
    slot_replacements_98.append((biggest, "模板选取决策树·场景→结构→推荐"))

# 5 BOLD labels → scenarios
for i, bs in enumerate(bold98[:5]):
    if i < len(decisions):
        name = decisions[i]['scenario']
        slot_replacements_98.append((bs, name))

# For each BOLD label, find a nearby desc slot and put path + recommendation
for i in range(min(5, len(decisions))):
    bs = bold98[i]
    nearest_desc = None
    min_dist = 999
    for ds in desc98:
        if any(ds is s for s, _ in slot_replacements_98):
            continue
        v_dist = abs(ds['top'] - bs['top'])
        if v_dist < 1.5 and v_dist < min_dist:
            min_dist = v_dist
            nearest_desc = ds
    
    if nearest_desc:
        desc_text = f"{decisions[i]['path']}·{decisions[i]['rec']}"
        slot_replacements_98.append((nearest_desc, desc_text))


# ============================================================
# Open PPT and keep only target slides
# ============================================================
prs = Presentation(TEMPLATE_PPT)
sldIdLst = prs.slides._sldIdLst
all_ids = list(sldIdLst)
rid163 = all_ids[PPT_IDX_163].get(qn('r:id'))
rid98 = all_ids[PPT_IDX_98].get(qn('r:id'))

for sld_id in reversed(all_ids):
    rid = sld_id.get(qn('r:id'))
    if rid != rid163 and rid != rid98:
        sldIdLst.remove(sld_id)

print(f"\nPPT: {len(list(sldIdLst))} slides retained")

# The slides are in original PPT order: S98 (idx 61) first, then S163 (idx 100)
slide_a = prs.slides[0]  # S98
slide_b = prs.slides[1]  # S163

# Verify
for i, s in enumerate([slide_a, slide_b]):
    for sh in s.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            print(f"  Slide {i}: \"{sh.text_frame.text[:50].strip()}\"")
            break

# ============================================================
# Apply replacements to S163 (slide_b)
# ============================================================
print("\n--- S163 replacements ---")
count163 = 0
for cat_slot, new_text in slot_replacements_163:
    shape = find_shape_by_pos_text(slide_b, cat_slot)
    if shape:
        old = shape.text_frame.paragraphs[0].text.strip()[:40]
        replace_shape_text(shape, new_text)
        count163 += 1
        print(f"  ✅ [{cat_slot['shape_name'][:20]:20s}] \"{old}\" → \"{new_text[:35]}\"")
    else:
        print(f"  ❌ [{cat_slot['shape_name'][:20]:20s}] NOT FOUND: \"{cat_slot['text'][:35]}\"")

print(f"  S163: {count163}/{len(slot_replacements_163)} replaced")

# ============================================================
# Apply replacements to S98 (slide_a)
# ============================================================
print("\n--- S98 replacements ---")
count98 = 0
for cat_slot, new_text in slot_replacements_98:
    shape = find_shape_by_pos_text(slide_a, cat_slot)
    if shape:
        old = shape.text_frame.paragraphs[0].text.strip()[:40]
        replace_shape_text(shape, new_text)
        count98 += 1
        print(f"  ✅ [{cat_slot['shape_name'][:20]:20s}] \"{old}\" → \"{new_text[:35]}\"")
    else:
        print(f"  ❌ [{cat_slot['shape_name'][:20]:20s}] NOT FOUND: \"{cat_slot['text'][:35]}\"")

print(f"  S98: {count98}/{len(slot_replacements_98)} replaced")

# ============================================================
# Save
# ============================================================
prs.save(OUTPUT_PPT)
print(f"\n✅ Saved: {OUTPUT_PPT}")
print(f"   Total: {count163 + count98} replacements across 2 slides")

#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
通用模板分析引擎。
对任意 PPTX 文件执行完整 DNA 提取 + 结构类型推断，输出模板目录 JSON。

用法：
  python analyze_template.py --template <path> [--output <path>]

默认华为模板：
  python analyze_template.py

自定义模板：
  python analyze_template.py --template D:/my_template.pptx --output references/custom_template_catalog.json
"""

import argparse, json, os, sys
from pptx import Presentation
from collections import Counter

# ── 工具函数 ──────────────────────────────────────────────────

def get_slide_dna(slide):
    """提取一页的完整DNA：每个文本框的位置/字号/字体/原文字/字数"""
    text_slots = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for pi, para in enumerate(shape.text_frame.paragraphs):
            text = para.text
            if not text.strip():
                continue

            # 字号和字体取第一个run
            size_pt = 0.0
            font_name = ""
            bold = False
            italic = False
            color = ""
            if para.runs:
                r = para.runs[0]
                size_pt = r.font.size / 12700 if r.font.size else 0
                font_name = r.font.name if r.font.name else ""
                bold = bool(r.font.bold)
                italic = bool(r.font.italic)
                try:
                    color = str(r.font.color.rgb) if r.font.color and r.font.color.rgb else ""
                except:
                    color = "SCHEME"

            left = shape.left / 914400 if shape.left else 0
            top = shape.top / 914400 if shape.top else 0
            w = shape.width / 914400 if shape.width else 0
            h = shape.height / 914400 if shape.height else 0

            text_slots.append({
                "shape_name": shape.name,
                "para_index": pi,
                "text": text.strip(),
                "char_count": len(text.strip()),
                "font_size": round(size_pt, 1),
                "font_name": font_name,
                "bold": bold,
                "italic": italic,
                "color": color,
                "left": round(left, 2),
                "top": round(top, 2),
                "width": round(w, 2),
                "height": round(h, 2),
            })
    return text_slots


def classify_structure_type(slots, slide_num):
    """
    根据文本槽位特征判断适合的内容结构。
    优先级链式判断，兜底"通用内容"。

    返回: (struct_type, confidence, description)
    """
    if not slots:
        return ("纯图页", 0.5, "无文本框")

    # 按字号分组
    big_title = [s for s in slots if s["font_size"] >= 24 and s["bold"]]
    medium_title = [s for s in slots if 16 <= s["font_size"] < 24]
    body_text = [s for s in slots if 10 <= s["font_size"] < 16]
    small_text = [s for s in slots if s["font_size"] < 10 and s["font_size"] > 0]

    total_chars = sum(s["char_count"] for s in slots)
    n_slots = len(slots)

    # ── 1. 封面/标题页 ──
    if len(big_title) == 1 and n_slots <= 3 and total_chars < 100:
        return ("通用内容", 0.9, "封面：大标题+≤3槽+<100字")

    # ── 2. 章节分隔 ──
    if len(big_title) == 1 and len(body_text) == 0 and total_chars < 60:
        return ("通用内容", 0.9, f"章节分隔：大标题+无正文+{total_chars}字")

    # ── 3. KPI 数据展示 ──
    kpi_like = [s for s in slots if s["font_size"] >= 24 and not s["bold"] and s["top"] >= 1.5]
    if 2 <= len(kpi_like) <= 6 and n_slots <= 12 and total_chars < 300:
        return ("KPI数据展示", 0.8, f"KPI数据：{len(kpi_like)}个大数字+{n_slots}槽+{total_chars}字")

    # ── 4. 总分结构 ──
    if len(big_title) == 1 and 3 <= len(medium_title) <= 8:
        return ("总分结构", 0.8, f"总分：1大标题+{len(medium_title)}个分项")

    # ── 5. 总分总结构 ──
    if len(big_title) >= 1 and len(body_text) >= 4 and total_chars > 200:
        # 尾部有小结文本
        summary_hint = [s for s in slots if s["font_size"] < 16 and len(s["text"]) < 30 and s["bold"]]
        if summary_hint or total_chars > 300:
            return ("总分总结构", 0.7, f"总分总：{len(big_title)}标题+{len(body_text)}段+{total_chars}字")

    # ── 6. 并列结构 ──
    if len(medium_title) >= 3:
        # 检查字号一致性
        sizes = [s["font_size"] for s in medium_title]
        if max(sizes) - min(sizes) <= 4.0:
            # 检查字数方差
            lens = [s["char_count"] for s in medium_title]
            avg_len = sum(lens) / len(lens)
            var = sum((l - avg_len)**2 for l in lens) / len(lens)
            if var < 30:
                return ("并列结构", 0.8, f"并列：{len(medium_title)}个同等级分项")

    # ── 7. 通用内容（兜底）──
    return ("通用内容", 0.5, f"通用：{n_slots}槽+{total_chars}字")


# ── 主分析流程 ──────────────────────────────────────────────────

def analyze_template(template_path, output_path):
    """完整分析流程"""

    if not os.path.exists(template_path):
        print(f"❌ 模板文件不存在: {template_path}")
        sys.exit(1)

    print(f"📂 加载模板: {template_path}")
    prs = Presentation(template_path)

    n_slides = len(prs.slides)
    slide_w = prs.slide_width / 914400
    slide_h = prs.slide_height / 914400
    print(f"   页数: {n_slides}  尺寸: {slide_w:.2f}\" × {slide_h:.2f}\"")

    catalog = []
    struct_types = []

    print(f"\n🔍 逐页分析中...")
    for i, slide in enumerate(prs.slides):
        slots = get_slide_dna(slide)
        struct_type, confidence, desc = classify_structure_type(slots, i + 1)
        struct_types.append(struct_type)

        total_chars = sum(s["char_count"] for s in slots)

        entry = {
            "slide": i + 1,
            "struct_type": struct_type,
            "confidence": confidence,
            "description": desc,
            "slot_count": len(slots),
            "total_chars": total_chars,
            "text_slots": slots,
        }
        catalog.append(entry)

        # 每页进度
        slot_desc = f"{len(slots)}槽/{total_chars}字"
        print(f"   页{i+1:3d}  [{struct_type:6s}] {slot_desc:12s}  {desc}")

    # ── 统计摘要 ────────────────────────────────────────────────

    type_counts = Counter(struct_types)
    print(f"\n{'='*60}")
    print(f"📊 分析完成")
    print(f"{'='*60}")
    print(f"  总页数:      {n_slides}")
    print(f"  幻灯片尺寸:  {slide_w:.2f}\" × {slide_h:.2f}\"")
    print(f"  结构类型分布:")

    for st, count in type_counts.most_common():
        pct = count / n_slides * 100
        bar = "█" * int(pct / 2)
        print(f"    {st:10s}  {count:3d}页 ({pct:4.1f}%)  {bar}")

    # 槽位统计
    slot_counts = [e["slot_count"] for e in catalog]
    print(f"\n  槽位统计:")
    print(f"    最少: {min(slot_counts)}  最多: {max(slot_counts)}  平均: {sum(slot_counts)/len(slot_counts):.1f}")

    char_counts = [e["total_chars"] for e in catalog]
    print(f"  字数统计:")
    print(f"    最少: {min(char_counts)}  最多: {max(char_counts)}  平均: {sum(char_counts)/len(char_counts):.1f}")

    # ── 保存 ────────────────────────────────────────────────────

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(catalog, f, ensure_ascii=False, indent=2)

    file_size = os.path.getsize(output_path)
    size_str = f"{file_size/1024:.0f}KB" if file_size < 1024*1024 else f"{file_size/1024/1024:.1f}MB"
    print(f"\n✅ 目录已保存: {output_path} ({size_str})")

    return catalog


# ── CLI 入口 ────────────────────────────────────────────────────

if __name__ == "__main__":
    _SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

    parser = argparse.ArgumentParser(
        description="通用模板分析引擎 — 提取 PPTX 全文DNA并输出结构目录 JSON"
    )
    parser.add_argument(
        "--template", "-t",
        default=os.path.join(_SKILL_DIR, "template", "hw_template.pptx"),
        help="模板 PPTX 路径（默认: template/hw_template.pptx）"
    )
    parser.add_argument(
        "--output", "-o",
        default=os.path.join(_SKILL_DIR, "references", "template_catalog_v2.json"),
        help="输出 JSON 路径（默认: references/template_catalog_v2.json）"
    )
    args = parser.parse_args()

    analyze_template(args.template, args.output)

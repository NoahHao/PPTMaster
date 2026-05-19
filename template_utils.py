#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""华为胶片模板工具集 — 共享函数"""

from pptx import Presentation
from pptx.oxml.ns import qn
import shutil, json

# ===== 文本替换（保留所有格式）=====

def replace_text_keep_font(para, new_text):
    """替换段落文本，完全保留字体/字号/颜色/粗斜体"""
    for run in para.runs:
        run.text = ""
    if para.runs:
        para.runs[0].text = new_text


# ===== 智能文本蒸馏（绝不用省略号）=====

def distill(target_len, text):
    """智能蒸馏到目标字数。核心约束：绝不在词语中间截断。
    - 短槽(≤8字)：整个词/短语，找自然断点
    - 中槽(9-20字)：找标点断句
    - 长槽(>20字)：找句号断句，允许115%超长"""
    if len(text) <= target_len:
        return text

    # === 极短槽 (≤8字)：取完整词 ===
    if target_len <= 8:
        # 找出 target_len 范围内的最后一个标点/空格位置
        cut = text[:target_len]
        best = target_len
        for sep in ["，","。","；","、"," ","】","）","：",":","·"]:
            idx = cut.rfind(sep)
            if idx > 0:
                best = min(best, idx)
        if best < target_len:
            return text[:best]
        # 没有标点，最少留个完整词
        # For Chinese, each char is a word; just take all
        return text[:target_len]

    # === 短槽 (9-20字)：找分句点 ===
    if target_len <= 20:
        cut = text[:target_len]
        for sep in ["。","；","，","、"]:
            idx = cut.rfind(sep)
            if idx >= target_len * 0.5:
                return text[:idx+1]
        # 没找到标点 → 往前找空白
        for sep in [" ","·"]:
            idx = cut.rfind(sep)
            if idx > 0:
                return text[:idx]
        return text[:target_len]

    # === 正文槽 (>20字)：句号断句，允许适度超长 ===
    if len(text) <= int(target_len * 1.15):
        return text.strip()
    
    cut = text[:int(target_len * 1.1)]
    for sep in ["。","；"]:
        idx = cut.rfind(sep)
        if idx >= target_len * 0.5:
            return text[:idx+1]
    return text[:target_len].strip()


# ===== 槽位角色分类 =====

def classify_slot_role(s, all_slots):
    """根据字号/粗体/字数/位置推断文本槽逻辑角色"""
    sz = s["font_size"]
    bld = s["bold"]
    chars = s["char_count"]
    top = s["top"]
    max_sz = max(x["font_size"] for x in all_slots)

    # 主标题：最大字号+顶部
    if sz >= 24 and top < 1.5:
        return "主标题"

    # 分节标题：大字号+粗体
    if sz >= 16 and bld:
        return "分项标题"

    # 数据值：大字号非粗非顶部
    if sz >= 24 and not bld and top >= 1.5:
        return "数据值"

    # 描述正文：中等字号+长文本 (>12字)
    if 12 <= sz < 20 and chars > 12:
        return "描述文"

    # 分类标签：中等字号+粗体+短文本
    if 12 <= sz < 16 and bld:
        return "分类标签"

    # 短标签：任何字号但字数少
    if 10 <= sz < 20 and chars <= 12:
        return "短标签"

    # 注释/来源：小字号
    if sz < 10:
        return "注释"

    return "正文"


# ===== PPT 操作 =====

def extract_single_slide(template_path, output_path, slide_num):
    """从模板提取单页幻灯片（0-based PPT index）
    
    通过复制目标页所有 shape 到新 PPT 来确保只含单页。
    注意：此方法保留所有 shape 的格式/位置/内容，但使用空白布局。
    """
    import copy as _copy
    from lxml import etree as _etree

    src_prs = Presentation(template_path)
    src_slide = src_prs.slides[slide_num]

    # Create new presentation with same dimensions
    dst_prs = Presentation()
    dst_prs.slide_width = src_prs.slide_width
    dst_prs.slide_height = src_prs.slide_height

    # Add blank slide
    blank_layout = dst_prs.slide_layouts[6]  # blank layout
    dst_slide = dst_prs.slides.add_slide(blank_layout)

    # Remove default blank shapes, then copy all shapes from source
    for sp in list(dst_slide.shapes):
        sp._element.getparent().remove(sp._element)

    for shape in src_slide.shapes:
        el = _copy.deepcopy(shape._element)
        dst_slide.shapes._spTree.append(el)

    dst_prs.save(output_path)
    return dst_prs


def load_template_catalog(catalog_path):
    """加载模板目录"""
    with open(catalog_path, "r", encoding="utf-8") as f:
        return json.load(f)


def get_slide_slots(catalog, slide_num):
    """获取指定页的文本槽位（已排序）"""
    entry = catalog[slide_num - 1]
    slots = entry["text_slots"]
    slots.sort(key=lambda s: (s["top"], s["left"]))
    return slots, entry


def match_shapes_to_plan(slide, plan):
    """将slide中的文本shapes按坐标匹配到plan槽位"""
    text_shapes = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        left = float(shape.left) / 914400 if shape.left else 0
        top = float(shape.top) / 914400 if shape.top else 0
        for para in shape.text_frame.paragraphs:
            if para.text.strip():
                text_shapes.append((round(top, 2), round(left, 2), para))

    text_shapes.sort(key=lambda x: (x[0], x[1]))
    return text_shapes


# ===== 智能内容池 =====

# 按字数分池 — 杜绝蒸馏截断
POOLS = {
    1: ["端", "核", "存", "算", "网", "安", "基", "智"],
    2: [
        "终端", "消息", "通道", "技能", "浏览", "文档",
        "调度", "编排", "存储", "画像", "引擎", "记忆",
        "安全", "合规", "插件", "工具", "协议", "路由",
        "缓存", "加密", "签名", "过滤", "拦截", "事件",
        "队列", "通知", "限流", "熔断", "降级", "重试",
        "幂等", "压缩", "解压", "回滚", "审计", "日志",
    ],
    3: [
        "多智体", "工具链", "持久化", "可扩展", "自进化",
        "去中心", "异步化", "声明式",
    ],
    5: [
        "插件化架构", "去中心调度", "声明式配置", "异步非阻塞",
        "多Agent框架",
    ],
    7: [
        "容器化一键部署", "多智能体调度层", "插件化能力市场",
        "记忆持久化引擎", "多通道统一接入",
    ],
    10: [
        "拥抱智能打造协作平台", "多智能体全栈解决方案", "开放生态插件化架构",
    ],
}

def pick_from_pools(target_len, idx):
    """从对应字数池取词，自动蒸馏"""
    for width in sorted(POOLS.keys()):
        if width >= target_len:
            pool = POOLS[width]
            item = pool[idx % len(pool)]
            if len(item) <= target_len:
                return distill(target_len, item)
            # 继续找更大的池
    # 兜底：2字池
    return distill(target_len, POOLS[2][idx % len(POOLS[2])])


# ===== 验证 =====

def verify_no_empty(slide):
    """验证slide没有空文本槽（原模就有文本的位置）"""
    empty = sum(1 for s in slide.shapes if s.has_text_frame
                for p in s.text_frame.paragraphs if p.text.strip() == "")
    return empty

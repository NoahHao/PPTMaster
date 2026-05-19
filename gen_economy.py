#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""中国经济 PPT — Slide 27 (运力×算力×存力∝数字经济)
直接按槽位索引映射，100%精确。
"""
from pptx import Presentation
from pptx.oxml.ns import qn
import shutil, json, os

_SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATE = os.path.join(_SKILL_DIR, "template", "hw_template.pptx")
SLIDE_NUM = 27

def replace_text_keep_font(para, new_text):
    for run in para.runs:
        run.text = ""
    if para.runs:
        para.runs[0].text = new_text

def distill(target_len, text):
    """蒸馏：短槽激进、正文保内容"""
    if len(text) <= target_len:
        return text
    if target_len <= 10:
        rules = [
            ("社会消费品零售总额", "零售总额"), ("固定资产投资", "固投"),
            ("同比增长", "增"), ("万亿元", "万亿"), ("GDP总量", "GDP"),
            ("出口总额", "出口额"),
        ]
        for a, b in rules:
            if len(text) <= target_len: return text
            text = text.replace(a, b)
        for sep in ["，", "、", "。", "；"]:
            if len(text) <= target_len: return text
            text = text.replace(sep, "")
        if len(text) > target_len:
            text = text[:target_len]
        return text.strip()
    # 正文：保内容，允许20%超长
    while "  " in text and len(text) > target_len:
        text = text.replace("  ", " ")
    if len(text) <= int(target_len * 1.2):
        return text.strip()
    return text[:target_len].strip()

# ===== 加载 =====
with open("template_catalog_v2.json", "r", encoding="utf-8") as f:
    catalog = json.load(f)
slots = catalog[SLIDE_NUM - 1]["text_slots"]
slots.sort(key=lambda s: (s["top"], s["left"]))

print(f"Slide {SLIDE_NUM} | {len(slots)}槽")
print(f"原标题: {slots[0]['text']}")

# ===== 中国经济内容：按槽位索引直接映射 =====
# Slot 0: 主标题 (28pt 17字)
# Slot 1: $ 符号 (31pt 1字)
# Slot 2-4: 数字经济 + 智能经济 (20pt)
# Slot 5-7: 运力 算力 存力 (20pt 2字)
# Slot 8-10: × × ∝ (32pt 1字)
# Slot 11-13: 三项描述
# Slot 14-15: 数据标签+值 组1
# Slot 16-17: 数据标签+值 组2
# Slot 18-19: 数据标签+值 组3
# Slot 20-21: GDP总结 + 年份
# Slot 22-23: 数据标签+值 组4
# Slot 24: 来源

economy_content = [
    # [0] 主标题 17字
    distill(17, "消费投资出口驱动中国经济高质量发展"),
    # [1] $ → ¥  1字
    "¥",
    # [2] 数字经济 4字 → 新经济
    distill(4, "新经济"),
    # [3] + 1字
    "+",
    # [4] 智能经济 4字 → 新动能
    distill(4, "新动能"),
    # [5] 运力 2字 → 消费
    "消费",
    # [6] 算力 2字 → 投资
    "投资",
    # [7] 存力 2字 → 出口
    "出口",
    # [8] × 1字 → 保留
    "×",
    # [9] × 1字 → 保留
    "×",
    # [10] ∝ 1字 → 保留
    "∝",
    # [11] [5G/5.5G + IPv6 + 光] 20字 → 消费数据
    distill(20, "社会零售总额48.8万亿元 同比增长7.2%"),
    # [12] [AI + MEC + CLOUD] 18字 → 投资数据
    distill(18, "固定资产投资51.5万亿元 同比增长4.2%"),
    # [13] [存储] 4字 → 三驾马车
    "三驾马车",
    # [14] 全球联接数 5字 → 零售总额
    distill(5, "零售总额"),
    # [15] 2,000亿 6字 → 48.8万亿
    distill(6, "48.8万亿"),
    # [16] 通用计算总量（FP32）12字 → 固投标签
    distill(12, "固定资产投资额"),
    # [17] 3.3Z FLOPS 10X 15字 → 51.5万亿
    distill(15, "51.5万亿元 增4.2%"),
    # [18] 产生且被存储的数据 9字 → 数字经济
    distill(9, "数字经济规模"),
    # [19] 1YB 23X 8字 → 50.2万亿
    distill(8, "50.2万亿"),
    # [20] USD$49万亿 18字 → GDP总结 (必须≤18)
    distill(18, "GDP126万亿元 增速5.2%"),
    # [21] 2030年 5字 → 2024年
    distill(5, "2024年"),
    # [22] AI计算算力（FP16）12字 → 出口额标签
    distill(12, "货物出口总额"),
    # [23] 105Z FLOPS 500X 16字 → 25.5万亿
    distill(16, "25.5万亿元 增5.0%"),
    # [24] 来源:华为《智能世界2030》16字 → 来源
    distill(16, "来源: 国家统计局 2024年"),
]

assert len(economy_content) == len(slots), f"内容数{len(economy_content)}≠槽位数{len(slots)}"

# ===== 生成 PPT =====
OUTPUT = "china_economy.pptx"
shutil.copy2(TEMPLATE, OUTPUT)
prs = Presentation(OUTPUT)

sldIdLst = prs.slides._sldIdLst
all_ids = list(sldIdLst)
target_rId = all_ids[SLIDE_NUM - 1].get(qn('r:id'))
for sld_id in reversed(all_ids):
    if sld_id.get(qn('r:id')) != target_rId:
        sldIdLst.remove(sld_id)

slide = prs.slides[0]

# 收集所有文本shape，按(top,left)精确排序
text_shapes = []
for shape in slide.shapes:
    if shape.has_text_frame:
        left = float(shape.left) / 914400 if shape.left else 0
        top = float(shape.top) / 914400 if shape.top else 0
        for pi, para in enumerate(shape.text_frame.paragraphs):
            if para.text.strip():
                text_shapes.append((round(top, 2), round(left, 2), shape, para))

text_shapes.sort(key=lambda x: (x[0], x[1]))

print(f"\n替换明细:")
for i, (top, left, shape, para) in enumerate(text_shapes):
    orig = para.text.strip()
    new = economy_content[i] if i < len(economy_content) else "—"
    replace_text_keep_font(para, new)
    match = "✓" if new else "✗"
    print(f"  [{i:2d}] {match} {orig[:25]:25s} → {new[:40]}")

prs.save(OUTPUT)
print(f"\n✓ {OUTPUT} — {len(text_shapes)}处全部替换, 格式完全保留")

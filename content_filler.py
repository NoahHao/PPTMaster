#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
内容填充器：将用户内容填入模板页面
策略：克隆模板slide → 识别文本占位区 → 替换文本
"""

from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
import copy
import re
from typing import List, Dict, Optional
from pathlib import Path


import os
_SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATE_PATH = os.path.join(_SKILL_DIR, "template", "hw_template.pptx")


class ContentFiller:
    """克隆模板页并填充用户内容"""

    TEXT_REPLACEMENT_RULES = {
        # 内容角色 → 对文本的处理方式
        "封面": "title_only",      # 只替换标题
        "总述": "title_body",      # 标题+正文
        "数据": "title_items",     # 标题+数据条目
        "分述": "title_body",
        "问题": "title_body",
        "方案": "title_body",
        "成效": "title_items",
        "总结": "title_only",
        "展望": "title_only",
        "过渡": "title_only",
        "背景": "title_only",
        "结尾": "title_only",
    }

    def __init__(self, template_path: str = None):
        self.template_path = template_path or TEMPLATE_PATH
        self.source_prs = Presentation(self.template_path)
        self.output_prs = Presentation(self.template_path)
        # 清空输出（保留第一页作为模板base）
        self._init_output()

    def _init_output(self):
        """初始化输出PPT，保留模板的slide master和layout"""
        # 删除所有现有slide
        xml_slides = self.output_prs.slides._sldIdLst
        if xml_slides is not None:
            for sld_id in list(xml_slides):
                xml_slides.remove(sld_id)

    def fill_slide(self, content_block: Dict, template_info: Dict) -> bool:
        """
        克隆一个模板页并填入内容

        Args:
            content_block: 内容块 {"role", "title", "body", "data_items", ...}
            template_info: 模板页信息 {"slide", "type", "layout", ...}

        Returns:
            bool: 是否成功
        """
        slide_num = template_info["slide"]
        role = content_block.get("role", "分述")

        # 克隆源slide
        source_slide = self.source_prs.slides[slide_num - 1]
        new_slide = self._clone_slide(source_slide)
        if new_slide is None:
            return False

        # 填充内容
        fill_mode = self.TEXT_REPLACEMENT_RULES.get(role, "title_body")
        self._apply_content(new_slide, content_block, fill_mode)

        return True

    def _clone_slide(self, source_slide):
        """克隆一个幻灯片到输出PPT"""
        # 获取源slide的layout
        src_layout = source_slide.slide_layout
        
        # 在输出PPT中创建新slide
        new_slide = self.output_prs.slides.add_slide(src_layout)
        
        # 复制所有shapes
        for shape in source_slide.shapes:
            el = copy.deepcopy(shape._element)
            new_slide.shapes._spTree.append(el)
        
        return new_slide

    def _apply_content(self, slide, content_block: Dict, fill_mode: str):
        """将内容填入slide的文本区域"""
        title = content_block.get("title", "")
        body = content_block.get("body", "")
        data_items = content_block.get("data_items", [])
        role = content_block.get("role", "分述")

        # 分析slide中可填充的文本区
        text_areas = self._get_text_areas(slide)
        
        if fill_mode == "title_only":
            self._fill_title(text_areas, title, role)
        
        elif fill_mode == "title_body":
            self._fill_title(text_areas, title, role)
            self._fill_body(text_areas, body)
        
        elif fill_mode == "title_items":
            self._fill_title(text_areas, title, role)
            if data_items:
                self._fill_data_items(text_areas, data_items)
            elif body:
                self._fill_body(text_areas, body)

    def _get_text_areas(self, slide) -> List[Dict]:
        """获取slide中所有可编辑的文本区域"""
        areas = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    
                    size = 0
                    bold = False
                    if para.runs:
                        size = para.runs[0].font.size/12700 if para.runs[0].font.size else 0
                        bold = para.runs[0].font.bold or False
                    
                    areas.append({
                        "shape": shape,
                        "paragraph": para,
                        "text": text,
                        "size": size,
                        "bold": bold,
                        "shape_name": shape.name,
                    })
        return areas

    def _fill_title(self, text_areas: List[Dict], new_title: str, role: str):
        """替换标题文本"""
        if not new_title:
            return
        
        # 找最大的文本（标题通常是最大的）
        title_areas = [a for a in text_areas if a["size"] >= 18 and a["bold"]]
        if not title_areas:
            # fallback: 找第一个文本
            title_areas = [text_areas[0]] if text_areas else []
        
        for area in title_areas:
            # 清除现有文本并设置新标题
            para = area["paragraph"]
            for run in para.runs:
                run.text = ""
            if para.runs:
                para.runs[0].text = new_title
            else:
                # 没有run就创建一个
                from pptx.oxml.ns import qn
                pPr = para._p.find(qn('a:pPr'))
                run_el = para._p.makeelement(qn('a:r'), {})
                t_el = run_el.makeelement(qn('a:t'), {})
                t_el.text = new_title
                run_el.append(t_el)
                para._p.append(run_el)
            break  # 只替换第一个

    def _fill_body(self, text_areas: List[Dict], body: str):
        """替换正文文本"""
        if not body:
            return
        
        # 找正文区（非标题的文本区）
        body_areas = [a for a in text_areas 
                      if a["size"] < 24 or not a["bold"]]
        
        if not body_areas:
            return
        
        # 将body分行填入
        lines = [l.strip() for l in body.split("\n") if l.strip()]
        
        for i, area in enumerate(body_areas):
            if i >= len(lines):
                break
            para = area["paragraph"]
            # 清除并设置
            for run in para.runs:
                run.text = ""
            if para.runs:
                para.runs[0].text = lines[i]
            else:
                pPr = para._p.find(qn('a:pPr'))
                run_el = para._p.makeelement(qn('a:r'), {})
                t_el = run_el.makeelement(qn('a:t'), {})
                t_el.text = lines[i]
                run_el.append(t_el)
                para._p.append(run_el)

        # 如果还有多行，追加到最后一个文本区
        if len(lines) > len(body_areas) and body_areas:
            last_area = body_areas[-1]
            last_para = last_area["paragraph"]
            remaining = "\n".join(lines[len(body_areas):])
            if last_para.runs:
                last_para.runs[0].text += "\n" + remaining

    def _fill_data_items(self, text_areas: List[Dict], data_items: List[Dict]):
        """填入数据条目（KPI格式）"""
        if not data_items:
            return
        
        # KPI型：大数字 + 小标签
        # 找最大font的区域放数字，次大的放标签
        sorted_areas = sorted(text_areas, key=lambda a: a["size"], reverse=True)
        
        for i, item in enumerate(data_items):
            label = item.get("label", "")
            value = item.get("value", "")
            
            # 数字区（大字号）
            if i * 2 < len(sorted_areas):
                area = sorted_areas[i * 2]
                para = area["paragraph"]
                for run in para.runs:
                    run.text = ""
                if para.runs:
                    para.runs[0].text = value
                else:
                    pPr = para._p.find(qn('a:pPr'))
                    run_el = para._p.makeelement(qn('a:r'), {})
                    t_el = run_el.makeelement(qn('a:t'), {})
                    t_el.text = value
                    run_el.append(t_el)
                    para._p.append(run_el)
            
            # 标签区（小字号）
            if i * 2 + 1 < len(sorted_areas):
                area = sorted_areas[i * 2 + 1]
                para = area["paragraph"]
                for run in para.runs:
                    run.text = ""
                if para.runs:
                    para.runs[0].text = label
                else:
                    pPr = para._p.find(qn('a:pPr'))
                    run_el = para._p.makeelement(qn('a:r'), {})
                    t_el = run_el.makeelement(qn('a:t'), {})
                    t_el.text = label
                    run_el.append(t_el)
                    para._p.append(run_el)

    def save(self, output_path: str):
        """保存输出PPT"""
        self.output_prs.save(output_path)
        print(f"已保存: {output_path}")


# ===== 测试 =====
if __name__ == "__main__":
    filler = ContentFiller()
    
    # 测试填充一页
    content = {
        "role": "分述",
        "title": "智能运维平台进展",
        "body": "完成第一阶段开发\n上线3个客户\n反馈良好"
    }
    
    template = {
        "slide": 11,
        "type": "内容-图文卡片",
        "layout": "1_空白"
    }
    
    success = filler.fill_slide(content, template)
    if success:
        filler.save("test_output.pptx")
        print("测试成功！")

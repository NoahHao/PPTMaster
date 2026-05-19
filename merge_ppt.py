#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPT合并输出器：将选中的模板页序列合并为最终PPT
"""

from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.oxml.ns import qn
import copy
import os
from pathlib import Path
from typing import List, Dict


_SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATE_PATH = os.path.join(_SKILL_DIR, "template", "hw_template.pptx")


class PPTMerger:
    """合并多个模板页到单个PPT文件"""

    def __init__(self, template_path: str = None):
        self.template_path = template_path or TEMPLATE_PATH
        self.source_prs = Presentation(self.template_path)
        self.slide_width = self.source_prs.slide_width
        self.slide_height = self.source_prs.slide_height

    def merge(self, slide_selections: List[Dict], output_path: str,
              include_cover: bool = True, include_toc: bool = False,
              include_end: bool = True) -> str:
        """
        合并选中的模板页

        Args:
            slide_selections: [{"content": {...}, "template": {...}}, ...]
            output_path: 输出文件路径
            include_cover: 是否包含封面
            include_toc: 是否包含目录
            include_end: 是否包含结尾页

        Returns:
            输出文件路径
        """
        # 创建输出PPT
        output_prs = Presentation(self.template_path)
        
        # 清空所有slide
        xml_slides = output_prs.slides._sldIdLst
        if xml_slides is not None:
            for sld_id in list(xml_slides):
                xml_slides.remove(sld_id)

        # === 1. 封面 (Slide 5) ===
        if include_cover:
            self._clone_slide_to(output_prs, 5)

        # === 2. 目录 (可选) ===
        if include_toc:
            # 简单方案：生成纯文本目录页
            slide = output_prs.slides.add_slide(output_prs.slide_layouts[6])  # 空白布局
            self._add_simple_slide(slide, "目录", 
                                   "请手动补充目录内容\n\n"
                                   + "\n".join(f"· {s['content'].get('title', '')}" 
                                               for s in slide_selections if s['content'].get('title')))

        # === 3. 正文页 ===
        for item in slide_selections:
            slide_num = item["template"]["slide"]
            self._clone_slide_to(output_prs, slide_num)

        # === 4. 结尾页 (Slide 253 - 最后一页) ===
        if include_end:
            self._clone_slide_to(output_prs, 253)

        # === 保存 ===
        output_prs.save(output_path)
        print(f"合并完成: {output_path} ({len(output_prs.slides)}页)")
        return output_path

    def _clone_slide_to(self, target_prs: Presentation, slide_num: int):
        """将源模板的指定页克隆到目标PPT"""
        source_slide = self.source_prs.slides[slide_num - 1]
        src_layout = source_slide.slide_layout

        # 在目标PPT中匹配相同名字的layout
        layout_name = src_layout.name
        target_layout = None
        for layout in target_prs.slide_layouts:
            if layout.name == layout_name:
                target_layout = layout
                break
        if target_layout is None:
            target_layout = target_prs.slide_layouts[0]

        new_slide = target_prs.slides.add_slide(target_layout)

        # 复制所有shapes
        for shape in source_slide.shapes:
            el = copy.deepcopy(shape._element)
            new_slide.shapes._spTree.append(el)

    def _add_simple_slide(self, slide, title_text: str, body_text: str = ""):
        """添加简单文本页（兜底用）"""
        from pptx.util import Inches, Pt
        
        # 标题
        left = Inches(1)
        top = Inches(0.5)
        width = Inches(11)
        height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = title_text
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.bold = True

        # 正文
        if body_text:
            left = Inches(1)
            top = Inches(1.8)
            width = Inches(11)
            height = Inches(5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = body_text
            tf.paragraphs[0].font.size = Pt(14)


class SmartMerger(PPTMerger):
    """智能合并器：克隆+填充内容"""

    def merge_with_content(self, slide_selections: List[Dict], output_path: str,
                           title: str = "", include_cover: bool = True) -> str:
        """
        合并模板页并填入用户内容
        """
        import io
        import zipfile
        from lxml import etree
        
        # 使用全新方式：读取模板的slide master和layout，创建干净的输出
        
        # 方法: 基于模板创建输出，然后清理
        output_prs = Presentation(self.template_path)
        
        # 清空所有现有slide（使用低级别方式以避免冲突）
        sldIdLst = output_prs.slides._sldIdLst
        if sldIdLst is not None:
            sld_ids = list(sldIdLst)
            for sld_id in sld_ids:
                sldIdLst.remove(sld_id)
        
        # 重置slide计数
        # 手动管理 - 添加slides时不会冲突

        # === 封面 ===
        if include_cover:
            self._clone_and_fill(output_prs, 5, {
                "role": "封面", "title": title, "body": "", "data_items": []
            })

        # === 正文 ===
        for item in slide_selections:
            slide_num = item["template"]["slide"]
            content = item["content"]
            self._clone_and_fill(output_prs, slide_num, content)

        # === 结尾 ===
        self._clone_slide_to(output_prs, 253)

        output_prs.save(output_path)
        print(f"智能合并完成: {output_path} ({len(output_prs.slides)}页)")
        return output_path

    def _clone_and_fill(self, target_prs, slide_num: int, content: Dict):
        """克隆并填充内容"""
        source_slide = self.source_prs.slides[slide_num - 1]
        src_layout = source_slide.slide_layout

        # 找匹配layout
        layout_name = src_layout.name
        target_layout = None
        for layout in target_prs.slide_layouts:
            if layout.name == layout_name:
                target_layout = layout
                break
        if target_layout is None:
            target_layout = target_prs.slide_layouts[0]

        new_slide = target_prs.slides.add_slide(target_layout)

        # 复制shapes
        for shape in source_slide.shapes:
            el = copy.deepcopy(shape._element)
            new_slide.shapes._spTree.append(el)

        # 填充文本
        title = content.get("title", "")
        body = content.get("body", "")
        role = content.get("role", "分述")
        data_items = content.get("data_items", [])

        self._fill_slide_text(new_slide, title, body, role, data_items)

    def _fill_slide_text(self, slide, title: str, body: str, role: str, data_items: List):
        """向slide填充文本"""
        # 收集文本区
        text_areas = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    size = 0
                    bold = False
                    if para.runs:
                        size = para.runs[0].font.size/12700 if para.runs[0].font.size else 0
                        bold = para.runs[0].font.bold or False
                    text_areas.append({
                        "shape": shape, "paragraph": para,
                        "text": text, "size": size, "bold": bold
                    })

        if not text_areas:
            return

        # 标题填充（找最大的文本区）
        if title:
            title_areas = sorted(text_areas, key=lambda a: a["size"], reverse=True)
            if title_areas:
                area = title_areas[0]
                self._replace_para_text(area["paragraph"], title)

        # 正文填充（找其余文本区）
        if body and role not in ("封面", "总结", "展望", "过渡"):
            non_title = [a for a in text_areas if a["size"] < 24 or not a["bold"]]
            lines = [l.strip() for l in body.split("\n") if l.strip()]
            for i, area in enumerate(non_title):
                if i >= len(lines):
                    break
                self._replace_para_text(area["paragraph"], lines[i])

        # 数据条目填充
        if data_items and len(data_items) >= 2:
            sorted_areas = sorted(text_areas, key=lambda a: a["size"], reverse=True)
            for i, item in enumerate(data_items[:len(sorted_areas)//2]):
                if i*2 < len(sorted_areas):
                    self._replace_para_text(sorted_areas[i*2]["paragraph"], 
                                           item.get("value", ""))
                if i*2+1 < len(sorted_areas):
                    self._replace_para_text(sorted_areas[i*2+1]["paragraph"],
                                           item.get("label", ""))

    def _replace_para_text(self, para, new_text: str):
        """替换段落的全部文本"""
        for run in para.runs:
            run.text = ""
        if para.runs:
            para.runs[0].text = new_text


# ===== 测试 =====
if __name__ == "__main__":
    merger = SmartMerger()
    
    test_selections = [
        {
            "content": {"role": "数据", "title": "Q1核心指标", 
                       "body": "营收5000万\n客户200家\n团队50人",
                       "data_items": [
                           {"label": "营收", "value": "5000万"},
                           {"label": "客户", "value": "200家"},
                           {"label": "团队", "value": "50人"}
                       ]},
            "template": {"slide": 9, "type": "内容-KPI/数据"}
        },
        {
            "content": {"role": "分述", "title": "智能运维平台",
                       "body": "完成第一阶段开发\n上线3个客户"},
            "template": {"slide": 11, "type": "内容-图文卡片"}
        },
        {
            "content": {"role": "总结", "title": "达成率85%",
                       "body": "核心指标超预期"},
            "template": {"slide": 7, "type": "内容-大图+标语"}
        },
    ]
    
    merger.merge_with_content(test_selections, "test_smart_merge.pptx", 
                              title="2026年Q1工作汇报")
